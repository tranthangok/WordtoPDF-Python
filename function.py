from fpdf import FPDF
from docx import Document
import openpyxl
from pptx import Presentation
import os

class Converter:
    @staticmethod
    def convert_docx_to_pdf(input_path, output_path, update_progress=None):
        pdf = FPDF()
        pdf.add_page()
    
        font_path = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")

        if os.path.exists(font_path):
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)

            doc = Document(input_path)
            for paragraph in doc.paragraphs:
                pdf.multi_cell(0, 10, paragraph.text)

            pdf.output(output_path)
            pdf.close()
        else:
            raise FileNotFoundError(f"TTF Font file not found: {font_path}")

    @staticmethod
    def convert_xlsx_to_pdf(input_path, output_path, update_progress=None):
        pdf = FPDF()
        pdf.add_page()

        font_path = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")

        if os.path.exists(font_path):
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)

            wb = openpyxl.load_workbook(input_path)
            sheet = wb.active

            for row in sheet.iter_rows():
                for cell in row:
                    pdf.cell(200, 10, str(cell.value), ln=1)

            pdf.output(output_path)
            pdf.close()
        else:
            raise FileNotFoundError(f"TTF Font file not found: {font_path}")

    @staticmethod
    def convert_pptx_to_pdf(input_path, output_path, update_progress=None):
        pdf = FPDF()
        pdf.add_page()

        font_path = os.path.join(os.path.dirname(__file__), "DejaVuSans.ttf")

        if os.path.exists(font_path):
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)

            prs = Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.multi_cell(0, 10, shape.text, align='L')

            pdf.output(output_path)
            pdf.close()
        else:
            raise FileNotFoundError(f"TTF Font file not found: {font_path}")

